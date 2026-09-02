# -*- coding: utf-8 -*-
# 路演 PPT 生成器：python tools/build_ppt.py → docs/ppt/OneCrew_路演PPT.pptx
# 9 页 · 16:9 · 深蓝主题 · 微软雅黑（a:ea XML 修正保证中文字体生效）
# 内容事实源：docs/08-v2路演讲稿.md（三金句/对比表/技术亮点/Q&A）
import os
import subprocess

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(ROOT, 'docs', 'ppt')
OUT = os.path.join(OUT_DIR, 'OneCrew_路演PPT.pptx')

NAVY = RGBColor(0x0E, 0x1A, 0x2F)   # 背景
PANEL = RGBColor(0x16, 0x26, 0x3D)  # 卡片
ACCENT = RGBColor(0x4F, 0x7B, 0xD9) # 主色
LIGHT = RGBColor(0xBF, 0xD4, 0xFF)  # 浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUT = RGBColor(0x8F, 0xA3, 0xC0)    # 灰蓝
GOOD = RGBColor(0x6F, 0xD9, 0x9B)   # 对勾绿
FONT = 'Microsoft YaHei'

IMG_COVER = os.path.join(ROOT, 'exports', '1788294529552_产品图集·1.png')
IMG_END = os.path.join(ROOT, 'exports', '1788297850396_产品图集·4.png')
IMG_P2 = os.path.join(ROOT, 'exports', '1788294529558_产品图集·2.png')
IMG_P3 = os.path.join(ROOT, 'exports', '1788294529562_产品图集·3.png')
VID_EVIDENCE = os.path.join(ROOT, 'exports', '1788294544399_15秒带货短视频.mp4')
SHOT_05 = os.path.join(ROOT, 'screenshots', '05-workflow-done-artifacts.png')
SHOT_06 = os.path.join(ROOT, 'screenshots', '06-run-timeline-done.png')
FRAME = os.path.join(ROOT, 'roadshow', 'work', 'ppt_frame.png')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_ea(run):
    """python-pptx 只设 latin 字体，中文需要补 a:ea typeface 才不回退宋体。"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def bg(slide, color=NAVY):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def txt(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: [(text, size, color, bold), ...] 每项一段。"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (t, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        set_ea(r)
    return tb


def card(slide, x, y, w, h, fill=PANEL, radius=0.10):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        s.adjustments[0] = radius
    except Exception:
        pass
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def header(slide, tag, title, sub=None):
    """页眉：左侧竖条 + 小标 + 大标题（+ 副标）。"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.52), Inches(0.09), Inches(0.92))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    lines = [(tag, 14, LIGHT, True), (title, 30, WHITE, True)]
    if sub:
        lines.append((sub, 14, MUT, False))
    txt(slide, 0.82, 0.42, 11.9, 1.35, lines)
    footer(slide)


def footer(slide):
    txt(slide, 9.6, 7.08, 3.3, 0.35, [('OneCrew · SynNovator Wave 3', 10, MUT, False)], align=PP_ALIGN.RIGHT)


def pic_h(slide, path, x, y, h):
    """按高度等比贴图（素材统一 1080x1920 或 16:9）。"""
    return slide.shapes.add_picture(path, Inches(x), Inches(y), height=Inches(h))


# ---------- P1 封面 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
pic_h(s, IMG_COVER, 13.333 - 4.22, 0, 7.5)  # 9:16 竖图贴右缘
veil = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.11), 0, Inches(4.22), prs.slide_height)
veil.fill.solid()
veil.fill.fore_color.rgb = NAVY
veil.fill.transparency = 0  # 不遮，留原图
s.shapes._spTree.remove(veil._element)  # 取消遮挡：直接删掉
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(2.02), Inches(1.4), Inches(0.12))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
txt(s, 0.88, 2.3, 7.6, 3.2, [
    ('OneCrew', 64, WHITE, True),
    ('一人出海 AIGC 内容官', 28, LIGHT, True),
    ('一个人的内容小队：你说一句产品，八位 AI 队友接力，', 16, MUT, False),
    ('183 秒拿到 13 个可发布工件——4 张 AI 产品图 + 1 条 15 秒竖屏带货成片。', 16, MUT, False),
])
txt(s, 0.88, 6.5, 8.0, 0.6, [('SynNovator AIGC 技术赛道 · Wave 3 Agents', 14, MUT, False)])

# ---------- P2 痛点 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '01 · 痛点', '一个人做跨境，内容是最大的时间黑洞', 'TikTok Shop 时代：内容 = 流量 = 订单，但一个人根本发不过来')
NUMS = [
    ('≥1 条 / 日', 'TikTok Shop 官方要求', '短视频日更是硬门槛'),
    ('6-10 小时 / 周', '60% 卖家被内容卡住', '选题→文案→图→片→日历 全手搓'),
    ('≈40% 首年被封', '78% 卖家搞不定合规', '广告法 / 文化禁忌 踩坑即限权'),
]
for i, (big, mid, sub) in enumerate(NUMS):
    x = 0.7 + i * 4.1
    card(s, x, 2.15, 3.8, 3.3)
    txt(s, x + 0.25, 2.5, 3.3, 2.7, [
        (big, 40, ACCENT if i != 2 else RGBColor(0xE8, 0x6A, 0x6A), True),
        (mid, 17, WHITE, True),
        (sub, 13, MUT, False),
    ])
txt(s, 0.7, 5.9, 12.0, 0.8, [('痛点不在「写文案」——在日更成片 + 每一条都不能踩合规红线。', 18, LIGHT, True)])

# ---------- P3 竞品对比 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '02 · 竞品实测', '8 家主流工具全测一遍：全链路 × 合规，行业空白', 'Predis · Simplified · Canva · InVideo · Fliki · CapCut · HeyGen · KreadoAI（详见 docs/01 §调研）')
ROWS = [
    ('能力', '8 家主流工具', 'OneCrew'),
    ('市场调研', '✗ 0/8', '✓ 40 秒出报告'),
    ('五平台文案', '△ 部分 / 模板', '✓ 8 位队友接力'),
    ('AI 出图', '△ 5/8', '✓ Seedream ×4'),
    ('AI 成片（旁白+镜头+字幕）', '△ 5/8 各自为战', '✓ 端到端 15s 成片'),
    ('内容日历', '△ 2/8', '✓ CSV 可直发'),
    ('广告法 / 文化合规体检', '✗ 0/8', '✓ 行业唯一'),
    ('人工确认点', '✗ 0/8', '✓ 确认点 = 编辑点'),
    ('单次全流程成本', '订阅 $19+/月', '≈ ¥1.2 / 次'),
]
gt = s.shapes.add_table(len(ROWS), 3, Inches(0.7), Inches(2.0), Inches(12.0), Inches(4.9)).table
gt.columns[0].width = Inches(4.6)
gt.columns[1].width = Inches(3.7)
gt.columns[2].width = Inches(3.7)
for r, row in enumerate(ROWS):
    gt.rows[r].height = Inches(0.52)
    for c, val in enumerate(row):
        cell = gt.cell(r, c)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.18)
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = ACCENT
            color, bold = WHITE, True
        else:
            cell.fill.fore_color.rgb = PANEL if r % 2 else RGBColor(0x1B, 0x2E, 0x4A)
            if c == 2:
                color, bold = (GOOD if val.startswith('✓') else LIGHT), True
            elif c == 1:
                color, bold = MUT, False
            else:
                color, bold = WHITE, False
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = val
        run.font.name = FONT
        run.font.size = Pt(14 if r == 0 else 13)
        run.font.bold = bold
        run.font.color.rgb = color
        set_ea(run)

# ---------- P4 流水线 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '03 · 产品', '八位 AI 队友接力：一条流水线交付可发布内容', '上下文接力——导演复用视觉工坊的图，合规审的是真实产出')
STEPS = [
    ('1', '市场快研'), ('2', '人工确认'), ('3', '品牌声音'), ('4', '五平台文案'),
    ('5', '提示词工坊'), ('6', '视觉工坊 ×4'), ('7', '短片导演 15s'), ('8', '内容日历'), ('9', '合规体检'),
]
for i, (num, name) in enumerate(STEPS):
    row, col = divmod(i, 5)
    x = 0.7 + col * 2.44
    y = 2.25 + row * 1.85
    hot = num == '2'
    card(s, x, y, 2.2, 1.5, fill=RGBColor(0x24, 0x3B, 0x60) if hot else PANEL)
    txt(s, x + 0.15, y + 0.18, 1.9, 1.2, [
        (num, 22, ACCENT if not hot else RGBColor(0xF5, 0xC2, 0x6B), True),
        (name, 15, WHITE, True),
        ('确认点=编辑点' if hot else '', 11, MUT, False),
    ])
txt(s, 0.7, 6.35, 12.0, 0.8, [('183 秒全流程 · 13 个可发布工件 · 每一步落盘留痕（data/runs 可复核）', 17, LIGHT, True)])

# ---------- P5 实证 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '04 · 实证', '真实运行留痕，不是 PPT 数字', 'data/runs/*.json 逐事件落盘 · 前端「成本」页签直接记账')
CARDS = [('183 秒', '全流程耗时'), ('13 个', '可发布工件'), ('9/9', '步骤完成'), ('¥1.32', '10 次运行合计')]
for i, (big, sub) in enumerate(CARDS):
    x = 0.7 + i * 3.1
    card(s, x, 1.95, 2.85, 1.25)
    txt(s, x + 0.2, 2.08, 2.5, 1.0, [(big, 26, ACCENT, True), (sub, 13, WHITE, False)])
s.shapes.add_picture(SHOT_06, Inches(0.7), Inches(3.5), width=Inches(5.85))
s.shapes.add_picture(SHOT_05, Inches(6.8), Inches(3.5), width=Inches(5.85))
txt(s, 0.7, 6.92, 12.0, 0.4, [
    ('运行时间线（9 步全绿）', 11, MUT, False), (' ', 4, MUT, False),
])

# ---------- P6 真成片 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '05 · 媒体工程', '真·成片：可直接发布的带货短视频', 'AI 配音 + 程序化运镜 + 硬字幕，1080×1920 MP4 当场可播')
if not os.path.exists(FRAME):
    os.makedirs(os.path.dirname(FRAME), exist_ok=True)
    subprocess.run([
        os.path.join(ROOT, 'node_modules', 'ffmpeg-static', 'ffmpeg.exe'), '-y',
        '-ss', '5', '-i', VID_EVIDENCE, '-frames:v', '1', FRAME,
    ], check=True, capture_output=True)
for j, img in enumerate([FRAME, IMG_P2, IMG_P3]):
    pic_h(s, img, 5.1 + j * 2.72, 2.0, 4.55)
txt(s, 0.7, 2.1, 4.2, 4.6, [
    ('15 秒竖屏带货片', 20, WHITE, True),
    ('· Edge TTS 多语言配音（中/英/日/韩）', 14, LIGHT, False),
    ('· Ken Burns 运镜 + xfade 转场', 14, LIGHT, False),
    ('· ASS 硬字幕（Noto CJK / 雅黑）', 14, LIGHT, False),
    ('· BGM 自动 ducking 压混', 14, LIGHT, False),
    ('成片 14 秒出片——', 15, WHITE, True),
    ('1.2 元成本下的正确工程取舍；', 13, MUT, False),
    ('接 Seedance i2v 即升级真视频片段，架构已预留。', 13, MUT, False),
])

# ---------- P7 技术亮点 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '06 · 技术亮点', '可上线产品与 demo 的分水岭', 'Agent 图谱执行器 · 上下文接力 · 媒体工程 · 可靠性设计')
TECH = [
    ('Agent 图谱执行器', '9 步状态机 + append-only 事件日志；每步落盘，评审可复核 data/runs/*.json。'),
    ('上下文接力', '每步工件摘要注入下游 context：短片导演直接复用 4 张图，合规体检审真实产出——拒绝「各步独立伪协作」。'),
    ('媒体工程', 'Seedream 4.0 生图 0.2 元/张 + Edge TTS 免费配音 + 自研 ffmpeg 合成层（运镜/xfade/ASS/ducking）。'),
    ('可靠性设计', '双重试 + 结构校验 + 确定性 mock 兜底；无生图 Key 自动降级「构图小样」并明示——全流程永不中断。'),
]
for i, (t, d) in enumerate(TECH):
    row, col = divmod(i, 2)
    x = 0.7 + col * 6.2
    y = 2.1 + row * 2.35
    card(s, x, y, 5.9, 2.05)
    txt(s, x + 0.3, y + 0.22, 5.3, 1.7, [(t, 18, ACCENT, True), (d, 13, WHITE, False)])

# ---------- P8 商业模式 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
header(s, '07 · 商业模式', '按次订阅：成本透明到分，毛利结构成立', '真实记账（LLM+生图+TTS），前端「成本」页签可见每一笔')
card(s, 0.7, 2.15, 3.8, 3.1)
txt(s, 1.0, 2.5, 3.3, 2.5, [('¥1.2 / 次', 38, ACCENT, True), ('单次全流程成本', 16, WHITE, True), ('LLM 0.24 + 真图 0.8 + TTS 0', 12, MUT, False)])
card(s, 4.8, 2.15, 3.8, 3.1)
txt(s, 5.1, 2.5, 3.3, 2.5, [('¥49 / 月', 38, WHITE, True), ('订阅 30 次 / 月', 16, WHITE, True), ('免费版 1 次/天 引流', 12, MUT, False)])
card(s, 8.9, 2.15, 3.8, 3.1)
txt(s, 9.2, 2.5, 3.3, 2.5, [('>95%', 38, GOOD, True), ('订阅毛利率', 16, WHITE, True), ('边际成本 ≈ 单次记账', 12, MUT, False)])
txt(s, 0.7, 5.75, 12.0, 1.0, [
    ('锚点：UGC 外包 $45-212/条 · 代运营 ¥5k-30k/月 → OneCrew 是 1/100 量级。', 18, LIGHT, True),
])

# ---------- P9 收尾 ----------
s = prs.slides.add_slide(BLANK)
bg(s)
pic_h(s, IMG_END, 13.333 - 4.22, 0, 7.5)
txt(s, 0.88, 2.2, 7.8, 3.6, [
    ('一个人 = 一支内容小队', 44, WHITE, True),
    ('别人还在写文案，我们的用户已经在发视频。', 20, LIGHT, False),
    ('每一分钱、每一步、每一个工件，全部可复核。', 16, MUT, False),
])
txt(s, 0.88, 6.5, 8.0, 0.6, [('SynNovator AIGC 技术赛道 · Wave 3 Agents · OneCrew', 14, MUT, False)])
footer(s)

os.makedirs(OUT_DIR, exist_ok=True)
prs.save(OUT)
print('saved:', OUT, len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
